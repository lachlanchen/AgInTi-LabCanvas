from __future__ import annotations

from copy import deepcopy
from datetime import datetime, timezone
import json
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
from typing import Any
from urllib.parse import quote
import zipfile


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MAX_GENERATED_ASSET_COVERAGE = 0.70
SUPPORTED_LAYOUTS = {
    "title",
    "section",
    "content",
    "two_column",
    "image_focus",
    "quote",
}
SUPPORTED_ELEMENT_TYPES = {"line", "shape", "table", "text"}
SUPPORTED_ELEMENT_SHAPES = {
    "chevron",
    "down_arrow",
    "hexagon",
    "oval",
    "rectangle",
    "right_arrow",
    "rounded_rectangle",
    "triangle",
}
GENERATED_KINDS = {"imagegen", "image_generation", "aginti", "generated"}
FORBIDDEN_GENERATED_ROLES = {
    "background",
    "full_slide",
    "full-slide",
    "slide",
    "slide_background",
}

THEMES: dict[str, dict[str, Any]] = {
    "bright_scientific": {
        "background": "#F7F9FC",
        "surface": "#FFFFFF",
        "text": "#17212B",
        "muted": "#53606D",
        "accent": "#008C83",
        "secondary": "#E7A928",
        "line": "#D8E0E8",
        "font_heading": "Aptos Display",
        "font_body": "Aptos",
    },
    "clean_editorial": {
        "background": "#FFFFFF",
        "surface": "#F4F6F8",
        "text": "#1D242B",
        "muted": "#626D78",
        "accent": "#C9413A",
        "secondary": "#167C80",
        "line": "#D9DEE3",
        "font_heading": "Aptos Display",
        "font_body": "Aptos",
    },
}


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds").replace("+00:00", "Z")


def slugify(value: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "-", str(value or "").strip()).strip("-").lower()
    return slug[:64] or "presentation"


def default_manifest(title: str, objective: str = "", theme: str = "bright_scientific") -> dict[str, Any]:
    if theme not in THEMES:
        raise ValueError(f"Unknown presentation theme: {theme}")
    objective_text = objective.strip() or "State the audience need, evidence, and recommended action."
    return {
        "schema_version": 1,
        "title": title.strip() or "Untitled presentation",
        "subtitle": "",
        "audience": "",
        "language": "match requester",
        "theme": theme,
        "theme_overrides": {},
        "objective": objective_text,
        "interaction": {
            "start_immediately": True,
            "progress_message": (
                "I have started the deck with a bright scientific theme. "
                "You can send audience, color, style, logo, or content changes while I work."
            ),
            "ask_only_if_materially_blocked": True,
        },
        "image_generation_policy": {
            "allowed_for_material_assets": True,
            "allow_text_when_reviewed": True,
            "forbid_full_slide_generation": True,
            "keep_slide_text_editable": True,
        },
        "slides": [
            {
                "id": "title",
                "layout": "title",
                "title": title.strip() or "Untitled presentation",
                "subtitle": objective_text,
                "bullets": [],
                "assets": [],
            },
            {
                "id": "context",
                "layout": "content",
                "title": "Context and question",
                "bullets": [
                    "What matters to the audience",
                    "What decision this presentation should support",
                ],
                "assets": [],
            },
            {
                "id": "evidence",
                "layout": "content",
                "title": "Evidence",
                "bullets": [
                    "Replace with traceable findings, data, or examples",
                    "Separate evidence from inference",
                ],
                "assets": [],
            },
            {
                "id": "next-steps",
                "layout": "content",
                "title": "Recommendation and next steps",
                "bullets": [
                    "Recommended action",
                    "Owner, timing, and validation signal",
                ],
                "assets": [],
            },
        ],
    }


def initialize_presentation_workspace(
    output_dir: str | Path,
    *,
    title: str,
    objective: str = "",
    theme: str = "bright_scientific",
) -> dict[str, Any]:
    root = Path(output_dir).expanduser().resolve()
    root.mkdir(parents=True, exist_ok=True)
    assets_dir = root / "assets"
    assets_dir.mkdir(exist_ok=True)
    manifest = default_manifest(title, objective, theme)
    manifest_path = root / "presentation.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    (assets_dir / "README.md").write_text(
        "# Presentation assets\n\n"
        "Keep each photo, generated illustration, plot, icon, or diagram as a separate file.\n"
        "For generated assets, preserve the prompt/provenance in `presentation.json`; never generate a complete slide image.\n",
        encoding="utf-8",
    )
    return {
        "ok": True,
        "project_dir": str(root),
        "manifest": str(manifest_path),
        "assets_dir": str(assets_dir),
        "progress_message": manifest["interaction"]["progress_message"],
    }


def load_manifest(path: str | Path) -> tuple[dict[str, Any], Path]:
    manifest_path = Path(path).expanduser().resolve()
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError("Presentation manifest must be a JSON object")
    return payload, manifest_path


def merged_theme(manifest: dict[str, Any]) -> dict[str, Any]:
    theme_name = str(manifest.get("theme") or "bright_scientific")
    if theme_name not in THEMES:
        raise ValueError(f"Unknown presentation theme: {theme_name}")
    theme = deepcopy(THEMES[theme_name])
    overrides = manifest.get("theme_overrides")
    if isinstance(overrides, dict):
        for key, value in overrides.items():
            if key in theme and str(value).strip():
                theme[key] = str(value).strip()
    return theme


def validate_manifest(manifest: dict[str, Any], base_dir: str | Path) -> dict[str, Any]:
    base = Path(base_dir).expanduser().resolve()
    errors: list[str] = []
    warnings: list[str] = []
    title = str(manifest.get("title") or "").strip()
    if not title:
        errors.append("Deck title is required")
    try:
        merged_theme(manifest)
    except ValueError as exc:
        errors.append(str(exc))

    slides = manifest.get("slides")
    if not isinstance(slides, list) or not slides:
        errors.append("At least one slide is required")
        slides = []
    if len(slides) > 80:
        errors.append("A single deck may contain at most 80 slides")

    slide_ids: set[str] = set()
    for index, slide in enumerate(slides, start=1):
        prefix = f"slide {index}"
        if not isinstance(slide, dict):
            errors.append(f"{prefix} must be an object")
            continue
        slide_id = str(slide.get("id") or f"slide-{index}").strip()
        if slide_id in slide_ids:
            errors.append(f"{prefix} duplicates id {slide_id!r}")
        slide_ids.add(slide_id)
        layout = str(slide.get("layout") or "content").strip()
        if layout not in SUPPORTED_LAYOUTS:
            errors.append(f"{prefix} has unsupported layout {layout!r}")
        if layout != "title" and not str(slide.get("title") or "").strip():
            warnings.append(f"{prefix} has no title")
        bullets = slide.get("bullets") or []
        if not isinstance(bullets, list):
            errors.append(f"{prefix} bullets must be a list")
        elif len(bullets) > 8:
            warnings.append(f"{prefix} has more than eight bullets; split it for readability")

        assets = slide.get("assets") or []
        if not isinstance(assets, list):
            errors.append(f"{prefix} assets must be a list")
            continue
        for asset_index, asset in enumerate(assets, start=1):
            asset_prefix = f"{prefix} asset {asset_index}"
            if not isinstance(asset, dict):
                errors.append(f"{asset_prefix} must be an object")
                continue
            raw_path = str(asset.get("path") or "").strip()
            if not raw_path:
                errors.append(f"{asset_prefix} needs a path")
                continue
            asset_path = _resolve_asset_path(base, raw_path)
            if not asset_path.is_file():
                errors.append(f"{asset_prefix} does not exist: {raw_path}")
            box = _asset_box(asset, layout)
            if box["w"] <= 0 or box["h"] <= 0:
                errors.append(f"{asset_prefix} box width and height must be positive")
            if box["x"] < 0 or box["y"] < 0:
                errors.append(f"{asset_prefix} box cannot start outside the slide")
            if box["x"] + box["w"] > SLIDE_WIDTH_IN or box["y"] + box["h"] > SLIDE_HEIGHT_IN:
                errors.append(f"{asset_prefix} box extends beyond the slide boundary")
            provenance = asset.get("provenance") if isinstance(asset.get("provenance"), dict) else {}
            source_kind = str(
                provenance.get("kind")
                or asset.get("source_kind")
                or ""
            ).strip().casefold()
            if source_kind not in GENERATED_KINDS:
                continue
            role = str(asset.get("role") or "supporting_visual").strip().casefold()
            if role in FORBIDDEN_GENERATED_ROLES:
                errors.append(f"{asset_prefix} cannot use generated imagery as {role!r}")
            coverage = (box["w"] * box["h"]) / (SLIDE_WIDTH_IN * SLIDE_HEIGHT_IN)
            if coverage > MAX_GENERATED_ASSET_COVERAGE:
                errors.append(
                    f"{asset_prefix} covers {coverage:.0%} of the slide; generated assets must remain supporting material"
                )
            prompt = str(provenance.get("prompt") or "").strip()
            prompt_path = str(provenance.get("prompt_path") or "").strip()
            if not prompt and not prompt_path:
                errors.append(f"{asset_prefix} needs a preserved generation prompt or prompt_path")
            if prompt_path and not _resolve_asset_path(base, prompt_path).is_file():
                errors.append(f"{asset_prefix} prompt_path does not exist: {prompt_path}")
            contains_text = bool(provenance.get("contains_text"))
            if contains_text:
                transcript = str(provenance.get("text_transcript") or "").strip()
                if not transcript:
                    errors.append(f"{asset_prefix} contains generated text but has no text_transcript")
                if provenance.get("text_reviewed") is not True:
                    errors.append(f"{asset_prefix} contains generated text but text_reviewed is not true")
                warnings.append(
                    f"{asset_prefix} contains reviewed bitmap text; keep essential wording as editable slide text too"
                )

        elements = slide.get("elements") or []
        if not isinstance(elements, list):
            errors.append(f"{prefix} elements must be a list")
            continue
        for element_index, element in enumerate(elements, start=1):
            element_prefix = f"{prefix} element {element_index}"
            if not isinstance(element, dict):
                errors.append(f"{element_prefix} must be an object")
                continue
            element_type = str(element.get("type") or "").strip().casefold()
            if element_type not in SUPPORTED_ELEMENT_TYPES:
                errors.append(
                    f"{element_prefix} has unsupported type {element_type!r}"
                )
                continue
            if element_type == "line":
                points = {
                    key: float(element.get(key, 0))
                    for key in ("x1", "y1", "x2", "y2")
                }
                if min(points.values()) < 0:
                    errors.append(f"{element_prefix} cannot start outside the slide")
                if max(points["x1"], points["x2"]) > SLIDE_WIDTH_IN:
                    errors.append(f"{element_prefix} extends beyond the slide width")
                if max(points["y1"], points["y2"]) > SLIDE_HEIGHT_IN:
                    errors.append(f"{element_prefix} extends beyond the slide height")
                continue
            box = _element_box(element)
            if box["w"] <= 0 or box["h"] <= 0:
                errors.append(f"{element_prefix} box width and height must be positive")
            if box["x"] < 0 or box["y"] < 0:
                errors.append(f"{element_prefix} box cannot start outside the slide")
            if box["x"] + box["w"] > SLIDE_WIDTH_IN or box["y"] + box["h"] > SLIDE_HEIGHT_IN:
                errors.append(f"{element_prefix} box extends beyond the slide boundary")
            if element_type == "shape":
                shape_kind = str(element.get("shape") or "rectangle").strip().casefold()
                if shape_kind not in SUPPORTED_ELEMENT_SHAPES:
                    errors.append(
                        f"{element_prefix} has unsupported shape {shape_kind!r}"
                    )
            if element_type == "table":
                rows = element.get("rows")
                if not isinstance(rows, list) or not rows:
                    errors.append(f"{element_prefix} table rows must be a non-empty list")
                elif not all(isinstance(row, list) and row for row in rows):
                    errors.append(f"{element_prefix} table rows must contain non-empty lists")
                elif len({len(row) for row in rows}) != 1:
                    errors.append(f"{element_prefix} table rows must have equal column counts")

    return {
        "ok": not errors,
        "errors": errors,
        "warnings": warnings,
        "slide_count": len(slides),
        "title": title,
    }


def build_presentation(
    manifest_path: str | Path,
    *,
    output_dir: str | Path | None = None,
    render: bool = False,
    timeout: float = 180,
) -> dict[str, Any]:
    manifest, source_path = load_manifest(manifest_path)
    validation = validate_manifest(manifest, source_path.parent)
    if not validation["ok"]:
        return {"ok": False, "validation": validation, "manifest": str(source_path)}
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise ValueError(
            "python-pptx is required. Install the project with `python -m pip install -e .`."
        ) from exc

    destination = (
        Path(output_dir).expanduser().resolve()
        if output_dir
        else (source_path.parent / "build").resolve()
    )
    destination.mkdir(parents=True, exist_ok=True)
    theme = merged_theme(manifest)
    deck_name = f"{slugify(str(manifest.get('title') or 'presentation'))}.pptx"
    pptx_path = destination / deck_name

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]
    slides = manifest.get("slides") or []
    for index, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_background(slide, theme, RGBColor)
        _render_slide(
            slide,
            spec,
            index=index,
            total=len(slides),
            base_dir=source_path.parent,
            theme=theme,
            api={
                "RGBColor": RGBColor,
                "MSO_SHAPE": MSO_SHAPE,
                "MSO_CONNECTOR": MSO_CONNECTOR,
                "MSO_ANCHOR": MSO_ANCHOR,
                "PP_ALIGN": PP_ALIGN,
                "Inches": Inches,
                "Pt": Pt,
            },
        )
    prs.save(pptx_path)

    package_check = validate_pptx_package(pptx_path, expected_slides=len(slides))
    render_result = render_presentation(pptx_path, destination, timeout=timeout) if render else {}
    audit = {
        "ok": package_check["ok"] and (not render or render_result.get("ok", False)),
        "built_at": utc_now(),
        "manifest": str(source_path),
        "pptx": str(pptx_path),
        "validation": validation,
        "package_check": package_check,
        "render": render_result,
        "image_generation_policy": manifest.get("image_generation_policy") or {},
    }
    audit_path = destination / "presentation-build.json"
    audit_path.write_text(
        json.dumps(audit, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return {**audit, "audit": str(audit_path)}


def validate_pptx_package(path: str | Path, *, expected_slides: int | None = None) -> dict[str, Any]:
    pptx_path = Path(path).expanduser().resolve()
    errors: list[str] = []
    slide_count = 0
    if not pptx_path.is_file() or pptx_path.stat().st_size < 1000:
        errors.append("PPTX file is missing or unexpectedly small")
    else:
        try:
            with zipfile.ZipFile(pptx_path) as archive:
                names = archive.namelist()
                if "[Content_Types].xml" not in names:
                    errors.append("PPTX package is missing [Content_Types].xml")
                slide_count = len(
                    [
                        name
                        for name in names
                        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
                    ]
                )
        except zipfile.BadZipFile:
            errors.append("PPTX output is not a readable ZIP package")
    if expected_slides is not None and slide_count != expected_slides:
        errors.append(f"Expected {expected_slides} slides, found {slide_count}")
    return {
        "ok": not errors,
        "errors": errors,
        "slide_count": slide_count,
        "size_bytes": pptx_path.stat().st_size if pptx_path.is_file() else 0,
    }


def render_presentation(path: str | Path, output_dir: str | Path, *, timeout: float = 180) -> dict[str, Any]:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        return {"ok": False, "error": "LibreOffice is not installed", "previews": []}
    pptx_path = Path(path).expanduser().resolve()
    destination = Path(output_dir).expanduser().resolve()
    destination.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="labcanvas-lo-") as profile:
        profile_uri = "file://" + quote(str(Path(profile).resolve()))
        completed = subprocess.run(
            [
                soffice,
                "--headless",
                f"-env:UserInstallation={profile_uri}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(destination),
                str(pptx_path),
            ],
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
        )
    pdf_path = destination / f"{pptx_path.stem}.pdf"
    if completed.returncode != 0 or not pdf_path.is_file():
        return {
            "ok": False,
            "error": (completed.stderr or completed.stdout or "LibreOffice conversion failed").strip(),
            "previews": [],
        }
    previews: list[str] = []
    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm:
        preview_dir = destination / "previews"
        preview_dir.mkdir(exist_ok=True)
        prefix = preview_dir / "slide"
        preview = subprocess.run(
            [pdftoppm, "-png", "-r", "110", str(pdf_path), str(prefix)],
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
        )
        if preview.returncode == 0:
            previews = [str(item) for item in sorted(preview_dir.glob("slide-*.png"))]
    return {"ok": True, "pdf": str(pdf_path), "previews": previews}


def _resolve_asset_path(base: Path, raw_path: str) -> Path:
    candidate = Path(raw_path).expanduser()
    if not candidate.is_absolute():
        candidate = base / candidate
    return candidate.resolve()


def _asset_box(asset: dict[str, Any], layout: str) -> dict[str, float]:
    raw = asset.get("box")
    if isinstance(raw, dict):
        return {
            "x": float(raw.get("x", 7.65)),
            "y": float(raw.get("y", 1.45)),
            "w": float(raw.get("w", 4.9)),
            "h": float(raw.get("h", 4.85)),
        }
    if layout == "image_focus":
        return {"x": 5.45, "y": 1.35, "w": 7.1, "h": 5.35}
    return {"x": 7.65, "y": 1.45, "w": 4.9, "h": 4.85}


def _element_box(element: dict[str, Any]) -> dict[str, float]:
    raw = element.get("box")
    if not isinstance(raw, dict):
        raw = element
    return {
        "x": float(raw.get("x", 0)),
        "y": float(raw.get("y", 0)),
        "w": float(raw.get("w", 1)),
        "h": float(raw.get("h", 1)),
    }


def _hex_rgb(value: str, rgb_color: Any) -> Any:
    cleaned = str(value or "#000000").strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", cleaned):
        raise ValueError(f"Invalid hex color: {value}")
    return rgb_color.from_string(cleaned.upper())


def _set_slide_background(slide: Any, theme: dict[str, Any], rgb_color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_rgb(theme["background"], rgb_color)


def _add_text_box(
    slide: Any,
    *,
    text: str,
    box: tuple[float, float, float, float],
    theme: dict[str, Any],
    api: dict[str, Any],
    font_size: float,
    color: str,
    bold: bool = False,
    align: Any | None = None,
    name: str = "",
) -> Any:
    shape = slide.shapes.add_textbox(
        api["Inches"](box[0]),
        api["Inches"](box[1]),
        api["Inches"](box[2]),
        api["Inches"](box[3]),
    )
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = api["MSO_ANCHOR"].TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align or api["PP_ALIGN"].LEFT
    run = paragraph.runs[0]
    run.font.name = theme["font_heading"] if bold else theme["font_body"]
    run.font.size = api["Pt"](font_size)
    run.font.bold = bold
    run.font.color.rgb = _hex_rgb(color, api["RGBColor"])
    return shape


def _add_bullets(
    slide: Any,
    *,
    bullets: list[Any],
    box: tuple[float, float, float, float],
    theme: dict[str, Any],
    api: dict[str, Any],
    name: str,
) -> Any:
    shape = slide.shapes.add_textbox(
        api["Inches"](box[0]),
        api["Inches"](box[1]),
        api["Inches"](box[2]),
        api["Inches"](box[3]),
    )
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = api["Inches"](0.08)
    frame.margin_right = api["Inches"](0.04)
    for index, raw in enumerate(bullets):
        if isinstance(raw, dict):
            text = str(raw.get("text") or "")
            level = max(0, min(2, int(raw.get("level") or 0)))
        else:
            text = str(raw)
            level = 0
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level
        paragraph.font.name = theme["font_body"]
        paragraph.font.size = api["Pt"](21 if level == 0 else 17)
        paragraph.font.color.rgb = _hex_rgb(theme["text"], api["RGBColor"])
        paragraph.space_after = api["Pt"](11 if level == 0 else 6)
    return shape


def _add_picture_contained(slide: Any, path: Path, box: dict[str, float], api: dict[str, Any], name: str) -> Any:
    try:
        from PIL import Image

        with Image.open(path) as image:
            width_px, height_px = image.size
    except Exception:
        width_px, height_px = (4, 3)
    scale = min(box["w"] / max(width_px, 1), box["h"] / max(height_px, 1))
    width = max(0.1, width_px * scale)
    height = max(0.1, height_px * scale)
    x = box["x"] + (box["w"] - width) / 2
    y = box["y"] + (box["h"] - height) / 2
    picture = slide.shapes.add_picture(
        str(path),
        api["Inches"](x),
        api["Inches"](y),
        width=api["Inches"](width),
        height=api["Inches"](height),
    )
    picture.name = name
    return picture


def _alignment(value: str, api: dict[str, Any]) -> Any:
    choices = {
        "center": api["PP_ALIGN"].CENTER,
        "left": api["PP_ALIGN"].LEFT,
        "right": api["PP_ALIGN"].RIGHT,
    }
    return choices.get(str(value or "left").casefold(), api["PP_ALIGN"].LEFT)


def _vertical_alignment(value: str, api: dict[str, Any]) -> Any:
    choices = {
        "bottom": api["MSO_ANCHOR"].BOTTOM,
        "middle": api["MSO_ANCHOR"].MIDDLE,
        "top": api["MSO_ANCHOR"].TOP,
    }
    return choices.get(str(value or "top").casefold(), api["MSO_ANCHOR"].TOP)


def _configure_element_text(
    frame: Any,
    *,
    text: str,
    element: dict[str, Any],
    theme: dict[str, Any],
    api: dict[str, Any],
) -> None:
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = _vertical_alignment(
        str(element.get("valign") or "middle"),
        api,
    )
    margin = float(element.get("margin", 0.08))
    frame.margin_left = api["Inches"](margin)
    frame.margin_right = api["Inches"](margin)
    frame.margin_top = api["Inches"](margin)
    frame.margin_bottom = api["Inches"](margin)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = _alignment(str(element.get("align") or "left"), api)
    paragraph.space_after = api["Pt"](0)
    for run in paragraph.runs:
        run.font.name = str(
            element.get("font_name")
            or (
                theme["font_heading"]
                if element.get("bold")
                else theme["font_body"]
            )
        )
        run.font.size = api["Pt"](float(element.get("font_size", 16)))
        run.font.bold = bool(element.get("bold"))
        run.font.color.rgb = _hex_rgb(
            str(element.get("text_color") or theme["text"]),
            api["RGBColor"],
        )


def _render_elements(
    slide: Any,
    elements: list[Any],
    *,
    theme: dict[str, Any],
    api: dict[str, Any],
) -> None:
    shape_map = {
        "chevron": api["MSO_SHAPE"].CHEVRON,
        "down_arrow": api["MSO_SHAPE"].DOWN_ARROW,
        "hexagon": api["MSO_SHAPE"].HEXAGON,
        "oval": api["MSO_SHAPE"].OVAL,
        "rectangle": api["MSO_SHAPE"].RECTANGLE,
        "right_arrow": api["MSO_SHAPE"].RIGHT_ARROW,
        "rounded_rectangle": api["MSO_SHAPE"].ROUNDED_RECTANGLE,
        "triangle": api["MSO_SHAPE"].ISOSCELES_TRIANGLE,
    }
    for index, element in enumerate(elements, start=1):
        if not isinstance(element, dict):
            continue
        element_type = str(element.get("type") or "").strip().casefold()
        name = str(element.get("name") or f"element-{index}")[:255]
        if element_type == "line":
            connector = slide.shapes.add_connector(
                api["MSO_CONNECTOR"].STRAIGHT,
                api["Inches"](float(element.get("x1", 0))),
                api["Inches"](float(element.get("y1", 0))),
                api["Inches"](float(element.get("x2", 0))),
                api["Inches"](float(element.get("y2", 0))),
            )
            connector.name = name
            connector.line.color.rgb = _hex_rgb(
                str(element.get("color") or theme["line"]),
                api["RGBColor"],
            )
            connector.line.width = api["Pt"](float(element.get("width", 1.5)))
            if str(element.get("dash") or "").casefold() == "dash":
                connector.line.dash_style = 4
            continue

        box = _element_box(element)
        if element_type == "text":
            shape = slide.shapes.add_textbox(
                api["Inches"](box["x"]),
                api["Inches"](box["y"]),
                api["Inches"](box["w"]),
                api["Inches"](box["h"]),
            )
            shape.name = name
            _configure_element_text(
                shape.text_frame,
                text=str(element.get("text") or ""),
                element=element,
                theme=theme,
                api=api,
            )
            continue

        if element_type == "shape":
            shape_kind = str(element.get("shape") or "rectangle").strip().casefold()
            shape = slide.shapes.add_shape(
                shape_map[shape_kind],
                api["Inches"](box["x"]),
                api["Inches"](box["y"]),
                api["Inches"](box["w"]),
                api["Inches"](box["h"]),
            )
            shape.name = name
            fill_color = str(element.get("fill") or theme["surface"])
            if fill_color.casefold() == "none":
                shape.fill.background()
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _hex_rgb(
                    fill_color,
                    api["RGBColor"],
                )
            line_color = str(element.get("line") or theme["line"])
            if line_color.casefold() == "none":
                shape.line.fill.background()
            else:
                shape.line.color.rgb = _hex_rgb(
                    line_color,
                    api["RGBColor"],
                )
                shape.line.width = api["Pt"](float(element.get("line_width", 1)))
            text = str(element.get("text") or "")
            if text:
                _configure_element_text(
                    shape.text_frame,
                    text=text,
                    element=element,
                    theme=theme,
                    api=api,
                )
            continue

        if element_type == "table":
            rows = element.get("rows") or []
            row_count = len(rows)
            column_count = len(rows[0])
            table_shape = slide.shapes.add_table(
                row_count,
                column_count,
                api["Inches"](box["x"]),
                api["Inches"](box["y"]),
                api["Inches"](box["w"]),
                api["Inches"](box["h"]),
            )
            table_shape.name = name
            table = table_shape.table
            column_widths = element.get("column_widths")
            if isinstance(column_widths, list) and len(column_widths) == column_count:
                total_width = sum(float(value) for value in column_widths)
                if total_width > 0:
                    for column, value in zip(table.columns, column_widths):
                        column.width = api["Inches"](
                            box["w"] * float(value) / total_width
                        )
            header_rows = max(0, int(element.get("header_rows", 1)))
            for row_index, values in enumerate(rows):
                for column_index, value in enumerate(values):
                    cell = table.cell(row_index, column_index)
                    is_header = row_index < header_rows
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_rgb(
                        str(
                            element.get("header_fill") or theme["accent"]
                            if is_header
                            else element.get("body_fill") or theme["surface"]
                        ),
                        api["RGBColor"],
                    )
                    cell_element = {
                        "align": (
                            element.get("header_align", "center")
                            if is_header
                            else element.get("body_align", "left")
                        ),
                        "bold": is_header,
                        "font_size": (
                            element.get("header_font_size", 13)
                            if is_header
                            else element.get("body_font_size", 12)
                        ),
                        "margin": element.get("cell_margin", 0.05),
                        "text_color": (
                            element.get("header_text_color", "#FFFFFF")
                            if is_header
                            else element.get("body_text_color", theme["text"])
                        ),
                        "valign": "middle",
                    }
                    _configure_element_text(
                        cell.text_frame,
                        text=str(value),
                        element=cell_element,
                        theme=theme,
                        api=api,
                    )


def _render_slide(
    slide: Any,
    spec: dict[str, Any],
    *,
    index: int,
    total: int,
    base_dir: Path,
    theme: dict[str, Any],
    api: dict[str, Any],
) -> None:
    layout = str(spec.get("layout") or "content")
    title = str(spec.get("title") or "")
    subtitle = str(spec.get("subtitle") or "")
    bullets = spec.get("bullets") if isinstance(spec.get("bullets"), list) else []
    assets = spec.get("assets") if isinstance(spec.get("assets"), list) else []
    elements = spec.get("elements") if isinstance(spec.get("elements"), list) else []

    if layout == "title":
        accent = slide.shapes.add_shape(
            api["MSO_SHAPE"].RECTANGLE,
            api["Inches"](0),
            api["Inches"](0),
            api["Inches"](0.22),
            api["Inches"](SLIDE_HEIGHT_IN),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = _hex_rgb(theme["accent"], api["RGBColor"])
        accent.line.fill.background()
        _add_text_box(
            slide,
            text=title,
            box=(0.85, 1.4, 11.6, 1.6),
            theme=theme,
            api=api,
            font_size=34,
            color=theme["text"],
            bold=True,
            name="title",
        )
        if subtitle:
            _add_text_box(
                slide,
                text=subtitle,
                box=(0.9, 3.2, 10.8, 1.1),
                theme=theme,
                api=api,
                font_size=19,
                color=theme["muted"],
                name="subtitle",
            )
    elif layout == "section":
        _add_text_box(
            slide,
            text=title,
            box=(0.85, 2.15, 11.6, 1.0),
            theme=theme,
            api=api,
            font_size=32,
            color=theme["text"],
            bold=True,
            name="section-title",
        )
        if subtitle:
            _add_text_box(
                slide,
                text=subtitle,
                box=(0.9, 3.35, 10.5, 0.9),
                theme=theme,
                api=api,
                font_size=19,
                color=theme["muted"],
                name="section-subtitle",
            )
    elif layout == "quote":
        _add_text_box(
            slide,
            text=title,
            box=(0.8, 0.55, 11.8, 0.65),
            theme=theme,
            api=api,
            font_size=27,
            color=theme["text"],
            bold=True,
            name="title",
        )
        _add_text_box(
            slide,
            text=str(spec.get("quote") or subtitle),
            box=(1.15, 1.75, 10.9, 3.2),
            theme=theme,
            api=api,
            font_size=27,
            color=theme["accent"],
            name="quote",
        )
        _add_text_box(
            slide,
            text=str(spec.get("attribution") or ""),
            box=(7.1, 5.2, 4.9, 0.6),
            theme=theme,
            api=api,
            font_size=16,
            color=theme["muted"],
            align=api["PP_ALIGN"].RIGHT,
            name="attribution",
        )
    else:
        _add_text_box(
            slide,
            text=title,
            box=(0.75, 0.42, 11.85, 0.7),
            theme=theme,
            api=api,
            font_size=27,
            color=theme["text"],
            bold=True,
            name="title",
        )
        line = slide.shapes.add_shape(
            api["MSO_SHAPE"].RECTANGLE,
            api["Inches"](0.77),
            api["Inches"](1.18),
            api["Inches"](1.15),
            api["Inches"](0.06),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_rgb(theme["accent"], api["RGBColor"])
        line.line.fill.background()
        if layout == "two_column":
            left = spec.get("left") if isinstance(spec.get("left"), list) else bullets
            right = spec.get("right") if isinstance(spec.get("right"), list) else []
            _add_bullets(
                slide,
                bullets=left,
                box=(0.8, 1.55, 5.7, 4.95),
                theme=theme,
                api=api,
                name="left-content",
            )
            _add_bullets(
                slide,
                bullets=right,
                box=(6.8, 1.55, 5.7, 4.95),
                theme=theme,
                api=api,
                name="right-content",
            )
        else:
            body_width = 4.35 if assets else 11.4
            _add_bullets(
                slide,
                bullets=bullets,
                box=(0.8, 1.55, body_width, 4.95),
                theme=theme,
                api=api,
                name="content",
            )

    for asset_index, asset in enumerate(assets, start=1):
        if not isinstance(asset, dict):
            continue
        path = _resolve_asset_path(base_dir, str(asset.get("path") or ""))
        box = _asset_box(asset, layout)
        _add_picture_contained(
            slide,
            path,
            box,
            api,
            name=f"asset-{asset_index}-{path.stem}"[:255],
        )
        caption = str(asset.get("caption") or "").strip()
        if caption:
            _add_text_box(
                slide,
                text=caption,
                box=(box["x"], min(6.55, box["y"] + box["h"] + 0.08), box["w"], 0.35),
                theme=theme,
                api=api,
                font_size=11,
                color=theme["muted"],
                name=f"asset-{asset_index}-caption",
            )

    _render_elements(slide, elements, theme=theme, api=api)

    _add_text_box(
        slide,
        text=f"{index} / {total}",
        box=(11.75, 7.02, 0.85, 0.25),
        theme=theme,
        api=api,
        font_size=9,
        color=theme["muted"],
        align=api["PP_ALIGN"].RIGHT,
        name="slide-number",
    )
